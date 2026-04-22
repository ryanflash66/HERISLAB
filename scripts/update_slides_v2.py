"""
In-place v2 metric update for HERISLAB_Progress_Update.pptx.

Skips slide 9 entirely (Matias's "Test Set Expansion" — must not be touched).
Backs up the original first.

Strategy: for each (slide_idx, old_text, new_text) rule, walk the slide's
shapes/text frames/runs and replace at the run level (preserving font,
color, size). Tables are updated by writing to cell text frames.
"""

import shutil
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPTX_PATH = ROOT / "docs" / "presentations" / "HERISLAB_Progress_Update.pptx"
BACKUP_PATH = ROOT / "docs" / "presentations" / "HERISLAB_Progress_Update.v1_backup.pptx"

# ---- Replacement rules ----
# Format: (slide_idx_1based, old_text, new_text)
# Slide 9 (Matias's "Test Set Expansion") is excluded everywhere.
TEXT_RULES = [
    # === Slide 4: Dataset Overview ===
    (4, "9,099", "9,117"),                                  # train stat card
    (4, "485", "722"),                                      # test stat card
    (4, "(33 Normal + 452 Fault)", "(37 Normal + 685 Fault)"),
    (4, "(33 Normal + 452 Fault)", "(37 Normal + 685 Fault)"),
    (4, "Preprocessing: Scale to 0-255  |  Resize to 320x240  |  Z-score normalize (mean=171.97, std=45.68)",
        "Preprocessing: Scale to 0-255  |  Resize to 320x240  |  Z-score normalize (mean=171.73, std=45.97)"),

    # === Slide 6: Training Results ===
    (6, "47", "50"),
    (6, "Early stopped", "Full 50 epochs"),
    (6, "(best at epoch 37)", "(best at epoch 50)"),
    (6, "0.0075", "0.00583"),
    (6, "epoch 37", "epoch 50"),
    (6, "Training converged smoothly. Loss dropped from 0.267 (epoch 1) to 0.0075 (epoch 37). Early stopping triggered at epoch 47 after 10 epochs with no improvement, confirming the model found its optimal point.",
        "Training converged smoothly across 50 epochs with a learning-rate drop to 5e-4 at epoch 33. Best validation loss reached 0.00583 at the final epoch, slightly improving on v1's 0.0075 despite a 50% larger fault test set covering a new equipment type."),

    # === Slide 7: Evaluation Results ===
    (7, "0.909", "0.902"),
    (7, "0.975", "0.981"),
    (7, "0.952", "0.962"),
    # Recall stays 1.000 — no rule needed
    (7, "Mean error: 0.001043", "Mean error: 0.000830"),
    (7, "Mean error: 0.004214  (4x higher)", "Mean error: 0.002810  (3.4x higher)"),

    # === Slide 10: Known Limitations (was old slide 9, shifted by Matias's insert) ===
    (10, "The autoencoder is currently trained on electric motors, induction motors, and photovoltaic panels. It learns what \"normal\" looks like for these specific equipment types. If presented with an unfamiliar equipment type (e.g., a transformer or circuit breaker), the model will likely flag it as anomalous even if the equipment is healthy, simply because it has never seen that visual pattern before.",
        "The autoencoder is currently trained on electric motors, induction motors, transformers, and photovoltaic panels. It learns what \"normal\" looks like for these specific equipment types. If presented with an unfamiliar equipment type (e.g., a circuit breaker or surge arrester), the model will likely flag it as anomalous even if the equipment is healthy, simply because it has never seen that visual pattern before."),
    (10, "Add normal images of new equipment types and retrain the model",
        "Add normal images of new equipment types and retrain. Transformers landed in v2; 879 more equipment images pending labels (HER-75)."),

    # === Slide 11: Team Contributions ===
    (11, "Dataset preparation support",
        "Image collection delivery: motors, transformers, 5 new equipment types"),

    # === Slide 12: Next Steps ===
    (12, "Generate reconstruction error heatmaps for visual explainability",
        "Heatmap generation (HER-25) - baseline shipped; refine with global error scaling"),
    (12, "Expand normal test set for more reliable metrics",
        "Expand test set with new equipment types (HER-75 pending dataset provenance)"),
]

# Table updates: (slide_idx, shape_idx, row_idx, col_idx, new_text)
# Slide 4's training source table needs: update counts + add transformer row.
TABLE_RULES = [
    # Slide 4 shape 16 — training source table
    (4, 16, 1, 0, "Electric Motor"),
    (4, 16, 1, 1, "168"),                  # was 176
    (4, 16, 1, 2, "PNG"),
    (4, 16, 2, 0, "Induction Motor (Noload)"),
    (4, 16, 2, 1, "20"),                   # was 25 — train portion only
    (4, 16, 2, 2, "BMP"),
    (4, 16, 3, 0, "Transformer (Noload)"), # NEW — overwrites old PV O&M row
    (4, 16, 3, 1, "18"),
    (4, 16, 3, 2, "BMP"),
    (4, 16, 4, 0, "PV O&M Inspection"),    # was PV Thermal Inspection — shift down
    (4, 16, 4, 1, "7,836"),
    (4, 16, 4, 2, "TIFF (16-bit)"),
]

# Slide 7 confusion-matrix table updates
TABLE_RULES.extend([
    (7, 16, 1, 1, "TP: 685"),  # was TP: 452
    (7, 16, 2, 1, "FP: 27"),   # was FP: 23
])

# Confusion matrix col 2 — FN/TN unchanged (FN: 0, TN: 10) so no rule

# Confirm we never accidentally edit slide 9
SKIP_SLIDE = 9


def replace_in_paragraph(para, old, new):
    """Replace `old` with `new` in a paragraph, preserving run formatting.

    Handles the common case where the entire match lives inside a single run.
    For multi-run matches (rare with the rules above), falls back to combining
    runs into the first run and clearing the rest.
    """
    # Single-run case
    for run in para.runs:
        if old in run.text:
            run.text = run.text.replace(old, new)
            return True

    # Multi-run case: stitch runs together, replace in stitched text
    full = "".join(r.text for r in para.runs)
    if old in full:
        new_full = full.replace(old, new)
        if para.runs:
            para.runs[0].text = new_full
            for r in para.runs[1:]:
                r.text = ""
        return True
    return False


def apply_text_rules(prs, rules):
    """Walk all shapes, apply text rules. Skip SKIP_SLIDE."""
    applied = 0
    for slide_idx_0, slide in enumerate(prs.slides):
        slide_idx = slide_idx_0 + 1
        if slide_idx == SKIP_SLIDE:
            continue
        slide_rules = [r for r in rules if r[0] == slide_idx]
        if not slide_rules:
            continue

        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for (_, old, new) in slide_rules:
                        if replace_in_paragraph(para, old, new):
                            applied += 1
            # Tables (text rules occasionally hit table cells too)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for (_, old, new) in slide_rules:
                                if replace_in_paragraph(para, old, new):
                                    applied += 1
    return applied


def apply_table_rules(prs, rules):
    """Targeted cell rewrites by (slide, shape_idx, row, col)."""
    applied = 0
    for slide_idx_0, slide in enumerate(prs.slides):
        slide_idx = slide_idx_0 + 1
        if slide_idx == SKIP_SLIDE:
            continue
        slide_rules = [r for r in rules if r[0] == slide_idx]
        if not slide_rules:
            continue

        shapes_list = list(slide.shapes)
        for (_, shape_idx, row_idx, col_idx, new_text) in slide_rules:
            if shape_idx >= len(shapes_list):
                print(f"  WARN slide {slide_idx} shape {shape_idx} not found")
                continue
            shape = shapes_list[shape_idx]
            if not shape.has_table:
                print(f"  WARN slide {slide_idx} shape {shape_idx} not a table")
                continue
            cell = shape.table.cell(row_idx, col_idx)
            # Set text on first run of first paragraph to preserve formatting
            tf = cell.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = new_text
                # Clear extra runs in first paragraph
                for r in tf.paragraphs[0].runs[1:]:
                    r.text = ""
                # Clear extra paragraphs
                for extra_para in tf.paragraphs[1:]:
                    for r in extra_para.runs:
                        r.text = ""
            else:
                cell.text = new_text
            applied += 1
    return applied


def main():
    if not PPTX_PATH.exists():
        print(f"ERROR: pptx not found at {PPTX_PATH}")
        return

    # Backup
    if not BACKUP_PATH.exists():
        shutil.copy2(PPTX_PATH, BACKUP_PATH)
        print(f"Backup created: {BACKUP_PATH.name}")
    else:
        print(f"Backup already exists: {BACKUP_PATH.name} (kept as-is)")

    prs = Presentation(PPTX_PATH)
    print(f"Loaded {len(prs.slides)} slides from {PPTX_PATH.name}")
    print(f"Skipping slide {SKIP_SLIDE} (Matias's Test Set Expansion)\n")

    print("Applying text rules...")
    text_applied = apply_text_rules(prs, TEXT_RULES)
    print(f"  {text_applied} text replacements applied")

    print("\nApplying table rules...")
    table_applied = apply_table_rules(prs, TABLE_RULES)
    print(f"  {table_applied} table cells updated")

    prs.save(PPTX_PATH)
    print(f"\nSaved updated pptx: {PPTX_PATH.name}")


if __name__ == "__main__":
    main()
