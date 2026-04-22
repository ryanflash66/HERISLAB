"""Build the 2026-04-22 status update deck.

Short standalone deck (8 slides), separate from the main v2 progress deck.
Inherits the v2 deck's slide master so fonts and theme colors match, then
rebuilds content slides from scratch.

Slides:
  1. Title
  2. Recap since v2
  3. Specialist Models per Equipment Type
  4. Ensemble Pipeline: ML + Rule-Based
  5. Manufacturer Temperature Thresholds
  6. Track C: Vision Transformers (Planned)
  7. Next Steps & Asks
  8. Thank You
"""

from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "docs" / "presentations" / "HERISLAB_Progress_Update.pptx"
OUT = ROOT / "docs" / "presentations" / "HERISLAB_Status_Update_2026-04-22.pptx"

# Theme colors (approximate v2 palette)
BLUE = RGBColor(0x1F, 0x4E, 0x79)        # primary
ACCENT = RGBColor(0x2E, 0x75, 0xB6)      # accent blue
DARK = RGBColor(0x1A, 0x1A, 0x1A)        # body text
MID = RGBColor(0x59, 0x59, 0x59)         # subtitle/muted
LIGHT = RGBColor(0xA6, 0xA6, 0xA6)       # light grey
GREEN = RGBColor(0x2E, 0x7D, 0x32)       # success
ORANGE = RGBColor(0xE6, 0x7E, 0x22)      # ask / warning
RED = RGBColor(0xC0, 0x39, 0x2B)         # critical
BG_CARD = RGBColor(0xF4, 0xF6, 0xF8)     # card background
BG_ASK = RGBColor(0xFF, 0xF4, 0xE6)      # ask callout background

SLIDE_W_IN = 10.0
SLIDE_H_IN = 5.625


def clear_slides(prs):
    """Remove all existing slides : both sldIdLst entries and the underlying parts."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get(qn("r:id"))
        # drop the underlying slide part so its XML file is released from the package
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)


def blank_slide(prs):
    """Add a blank slide using the default layout."""
    return prs.slides.add_slide(prs.slide_layouts[0])


def add_textbox(slide, left, top, width, height, text, *,
                font_size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    """Add a text box with formatted text. text can be str or list of (text, size, bold, color) tuples."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = [text]
    else:
        lines = text

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align

        if isinstance(line, tuple):
            txt, sz, bd, clr = line
        else:
            txt, sz, bd, clr = line, font_size, bold, color

        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(sz)
        run.font.bold = bd
        if clr is not None:
            run.font.color.rgb = clr
    return tb


def add_bullets(slide, left, top, width, height, items, *, font_size=12, color=DARK,
                bullet_color=ACCENT):
    """Add a bulleted list. items is a list of strings or (str, color_override) tuples."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            txt, clr = item
        else:
            txt, clr = item, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.color.rgb = bullet_color
        bullet.font.size = Pt(font_size)
        bullet.font.bold = True
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.color.rgb = clr
        run.font.name = "Calibri"
    return tb


def add_rect(slide, left, top, width, height, fill=BG_CARD, line=None):
    """Add a rounded rectangle (card background)."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_title(slide, text, color=BLUE, size=28):
    """Standard slide title, top-left."""
    return add_textbox(slide, 0.4, 0.25, 9.2, 0.5, text,
                       font_size=size, bold=True, color=color)


def add_subtitle(slide, text, color=MID, size=13):
    """Slide subtitle, under the title."""
    return add_textbox(slide, 0.4, 0.8, 9.2, 0.35, text,
                       font_size=size, bold=False, color=color)


# ---- Slide builders ----

def slide_title(prs):
    s = blank_slide(prs)
    # Background accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0),
                             Inches(SLIDE_W_IN), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(s, 0.5, 1.0, 9.0, 0.9, "Status Update",
                font_size=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 1.7, 9.0, 0.5, "HERIS Lab  |  Thermal Fault Detection",
                font_size=20, bold=False, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 2.4, 9.0, 0.4, "April 22, 2026",
                font_size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 4.6, 9.0, 0.4,
                "Ryan Balungeli  |  Matias  |  Pablo",
                font_size=14, color=MID, align=PP_ALIGN.CENTER)
    return s


def slide_recap(prs):
    s = blank_slide(prs)
    add_title(s, "Recap Since v2")
    add_subtitle(s, "Current baseline + what's moved forward")

    # Left card: v2 metrics (still current)
    add_rect(s, 0.4, 1.3, 4.5, 3.7, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Current Baseline (v2)",
                font_size=16, bold=True, color=BLUE)
    add_textbox(s, 0.6, 1.85, 4.1, 0.3,
                "Trained on 4 equipment types, 9,117 normal images",
                font_size=10, color=MID)

    metrics = [
        ("AUROC", "0.902", ACCENT),
        ("F1 Score", "0.981", GREEN),
        ("Precision", "0.962", GREEN),
        ("Recall", "1.000", GREEN),
    ]
    y = 2.3
    for label, val, c in metrics:
        add_textbox(s, 0.7, y, 2.0, 0.35, label, font_size=12, color=DARK)
        add_textbox(s, 2.7, y, 2.0, 0.35, val, font_size=14, bold=True, color=c)
        y += 0.45

    add_textbox(s, 0.6, 4.3, 4.1, 0.25,
                "685 faults detected, 0 missed, 37 normals tested",
                font_size=9, color=MID)
    add_textbox(s, 0.6, 4.55, 4.1, 0.25,
                "Test set: 37 normal + 685 fault across motors / transformer / PV",
                font_size=9, color=MID)

    # Right: what's new (bullets)
    add_textbox(s, 5.1, 1.3, 4.5, 0.4, "What's Moved Forward",
                font_size=16, bold=True, color=BLUE)
    add_bullets(s, 5.1, 1.75, 4.5, 3.3, [
        "Equipment scope narrowed to transformer + PV only",
        "Transformer thresholds characterized (Matias, HI-105)",
        "PV thresholds characterized (first-pass, public sources)",
        "New architecture: ensemble ML + rule-based inference",
        "Track C scoped (ViT / PatchCore) as parallel baseline",
        "Reconstruction-error heatmap baseline shipped",
        "Known blockers identified with clear ownership",
    ], font_size=12)
    return s


def slide_specialist_models(prs):
    s = blank_slide(prs)
    add_title(s, "Specialist Models per Equipment Type")
    add_subtitle(s, "Why one unified model won't work for the narrowed scope")

    # Problem
    add_rect(s, 0.4, 1.3, 4.5, 1.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "The Problem",
                font_size=15, bold=True, color=RED)
    add_bullets(s, 0.6, 1.85, 4.1, 1.3, [
        "~8,898 PV normal images vs. ~18 transformer normals",
        "494:1 class imbalance",
        "A unified model would learn 'normal = PV' and flag transformers as anomalous regardless of health",
    ], font_size=11)

    # Solution
    add_rect(s, 5.1, 1.3, 4.5, 1.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "The Solution",
                font_size=15, bold=True, color=GREEN)
    add_bullets(s, 5.3, 1.85, 4.1, 1.3, [
        "Two separate autoencoders, one per equipment type",
        "Each sees only its own normal images during training",
        "Routing logic at inference (manual for v3, automated later)",
        "Matches new equipment scope + the 'Per-Equipment Models' mitigation strategy",
    ], font_size=11)

    # Blockers bar
    add_rect(s, 0.4, 3.4, 9.2, 1.8, fill=BG_ASK)
    add_textbox(s, 0.6, 3.55, 9.0, 0.4, "Blocked On",
                font_size=15, bold=True, color=ORANGE)
    add_bullets(s, 0.6, 3.95, 9.0, 1.15, [
        "More transformer data : Matias sourcing additional public datasets",
        "Augmentation pipeline : critical when training set is ~18-100 images",
        "Both are independently unblocked; augmentation can start tonight",
    ], font_size=11)
    return s


def slide_ensemble(prs):
    s = blank_slide(prs)
    add_title(s, "Ensemble Pipeline: ML + Rule-Based")
    add_subtitle(s, "Complementary, not redundant : they catch different failure modes")

    # Architecture boxes (horizontal flow)
    box_w = 1.8
    box_h = 0.9
    box_y = 1.5

    # Input
    add_rect(s, 0.3, box_y, box_w, box_h, fill=ACCENT)
    add_textbox(s, 0.3, box_y + 0.25, box_w, 0.4, "Raw Input",
                font_size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, box_y + 0.55, box_w, 0.3, "Thermal Image",
                font_size=9, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(s, 2.15, box_y + 0.25, 0.4, 0.4, "→",
                font_size=24, bold=True, color=MID, align=PP_ALIGN.CENTER)

    # ML path
    add_rect(s, 2.55, 1.0, 2.2, 0.6, fill=GREEN)
    add_textbox(s, 2.55, 1.1, 2.2, 0.4, "ML: Autoencoder",
                font_size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 2.6, 1.7, 2.1, 0.3, "Reconstruction error",
                font_size=10, color=DARK, align=PP_ALIGN.CENTER)

    # Rules path
    add_rect(s, 2.55, 2.4, 2.2, 0.6, fill=ORANGE)
    add_textbox(s, 2.55, 2.5, 2.2, 0.4, "Rules: Thresholds",
                font_size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 2.6, 3.1, 2.1, 0.3, "Max temp vs. manual limits",
                font_size=10, color=DARK, align=PP_ALIGN.CENTER)

    # Connector arrows
    add_textbox(s, 4.8, 1.1, 0.3, 0.4, "→", font_size=20, bold=True, color=MID)
    add_textbox(s, 4.8, 2.5, 0.3, 0.4, "→", font_size=20, bold=True, color=MID)

    # Ensemble
    add_rect(s, 5.2, 1.5, 1.8, box_h, fill=BLUE)
    add_textbox(s, 5.2, 1.7, 1.8, 0.4, "Ensemble",
                font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 5.2, 2.05, 1.8, 0.3, "Combined decision",
                font_size=9, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Arrow → output
    add_textbox(s, 7.05, 1.7, 0.3, 0.4, "→", font_size=20, bold=True, color=MID)

    # Final output
    add_rect(s, 7.4, 1.5, 2.2, box_h, fill=BG_CARD, line=DARK)
    add_textbox(s, 7.4, 1.65, 2.2, 0.35, "Output",
                font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, 7.4, 2.0, 2.2, 0.45, "Flag + Severity\n+ Explanation",
                font_size=10, color=MID, align=PP_ALIGN.CENTER)

    # Bottom: key insight
    add_rect(s, 0.4, 3.7, 9.2, 1.6, fill=BG_CARD)
    add_textbox(s, 0.6, 3.85, 9.0, 0.4, "Why both layers matter",
                font_size=14, bold=True, color=BLUE)

    # Two columns below
    add_textbox(s, 0.6, 4.25, 4.4, 0.3, "ML catches:", font_size=11, bold=True, color=GREEN)
    add_bullets(s, 0.6, 4.5, 4.4, 0.7, [
        "Novel failure modes nobody wrote a rule for",
        "Localized hot spots below overall temp ceilings",
        "Asymmetric / spatial pattern anomalies",
    ], font_size=9)

    add_textbox(s, 5.1, 4.25, 4.4, 0.3, "Rules catch:", font_size=11, bold=True, color=ORANGE)
    add_bullets(s, 5.1, 4.5, 4.4, 0.7, [
        "Known failure modes with defined limits",
        "Auditable, regulator-friendly decisions",
        "Safety-critical hard ceilings (120°C oil, 125°C arrester)",
    ], font_size=9)
    return s


def slide_thresholds(prs):
    s = blank_slide(prs)
    add_title(s, "Manufacturer Temperature Thresholds")
    add_subtitle(s, "HI-105 transformer manual (Matias) + PV module research (first-pass)")

    # Transformer column
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Transformer (HI-105)",
                font_size=15, bold=True, color=BLUE)
    add_textbox(s, 0.6, 1.8, 4.1, 0.3,
                "Howard Industries three-phase padmount",
                font_size=9, color=MID)

    # Simple threshold rows
    rows = [
        ("Top-oil ceiling", "ambient + 85°C"),
        ("Summer estimate (35°C amb)", "120°C"),
        ("MOV arrester (recommended)", "90°C avg oil"),
        ("MOV arrester (max)", "125°C"),
        ("Min fluid (mineral oil)", "−20°C"),
        ("Min fluid (R-temp / FR3)", "0 / −10°C"),
    ]
    y = 2.2
    for label, val in rows:
        add_textbox(s, 0.7, y, 2.8, 0.28, label, font_size=10, color=DARK)
        add_textbox(s, 3.5, y, 1.3, 0.28, val, font_size=10, bold=True, color=BLUE)
        y += 0.3

    add_textbox(s, 0.6, 4.1, 4.1, 0.9,
                "Key insight: IR reads tank surface, not oil. Surface ≈ top-oil − 10 to 15°C under normal load. Ensemble pipeline uses +10°C calibration offset (conservative).",
                font_size=9, color=ORANGE)

    # PV column
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "PV Modules (Research-Compiled)",
                font_size=15, bold=True, color=BLUE)
    add_textbox(s, 5.3, 1.8, 4.1, 0.3,
                "Uniform across LG, JinkoSolar, Canadian Solar, Trina, First Solar",
                font_size=9, color=MID)

    # 4-tier table header
    add_textbox(s, 5.3, 2.2, 1.6, 0.28, "Component", font_size=9, bold=True, color=DARK)
    add_textbox(s, 6.9, 2.2, 0.7, 0.28, "Normal", font_size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, 7.6, 2.2, 0.7, 0.28, "Warn", font_size=9, bold=True, color=RGBColor(0xE6, 0xB0, 0x00), align=PP_ALIGN.CENTER)
    add_textbox(s, 8.3, 2.2, 0.7, 0.28, "Alarm", font_size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(s, 9.0, 2.2, 0.5, 0.28, "Crit", font_size=9, bold=True, color=RED, align=PP_ALIGN.CENTER)

    pv_rows = [
        ("Cell temp",   "<60°C",  "60–70",  "70–85",  ">85°C"),
        ("Bypass diode","<120°C", "120–150","150–175",">175°C"),
        ("Hot-spot ΔT", "<15°C",  "15–40",  "40–75",  ">75°C"),
    ]
    y = 2.55
    for comp, n, w, a, c in pv_rows:
        add_textbox(s, 5.3, y, 1.6, 0.28, comp, font_size=9, color=DARK)
        add_textbox(s, 6.9, y, 0.7, 0.28, n, font_size=9, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, 7.6, y, 0.7, 0.28, w, font_size=9, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, 8.3, y, 0.7, 0.28, a, font_size=9, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(s, 9.0, y, 0.5, 0.28, c, font_size=9, color=DARK, align=PP_ALIGN.CENTER)
        y += 0.32

    add_textbox(s, 5.3, 3.7, 4.1, 1.2,
                "Sources: IEC 61215, IEC 61730, UL 61730, 5 manufacturer datasheets, 7+ peer-reviewed papers. Per-model variation negligible for fault detection : single threshold set applies across mono-Si, poly-Si, thin-film.",
                font_size=9, color=MID)
    return s


def slide_track_c(prs):
    s = blank_slide(prs)
    add_title(s, "Track C: Vision Transformers (Planned)")
    add_subtitle(s, "Parallel baseline for head-to-head comparison against autoencoder")

    # Track C v1 column
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Track C v1",
                font_size=15, bold=True, color=GREEN)
    add_textbox(s, 0.6, 1.8, 4.1, 0.3,
                "Frozen DINOv2 + PatchCore",
                font_size=11, bold=True, color=DARK)
    add_bullets(s, 0.6, 2.2, 4.1, 2.5, [
        "Pretrained ViT (Meta, 142M images)",
        "No training : inference only on existing RTX 5070",
        "Memory-bank anomaly detection: stores normal patch features, flags distant new patches",
        "Per-patch heatmaps for free : cleaner than v2 per-image normalization",
        "SOTA on industrial anomaly benchmarks (MVTec AD)",
        "2–4 days effort, no new hardware",
    ], font_size=10)
    add_textbox(s, 0.6, 4.75, 4.1, 0.35, "Ready to start any time",
                font_size=10, bold=True, color=GREEN)

    # Track C v2 column
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "Track C v2",
                font_size=15, bold=True, color=ORANGE)
    add_textbox(s, 5.3, 1.8, 4.1, 0.3,
                "MAE / AnoDDPM / Reverse Distillation",
                font_size=11, bold=True, color=DARK)
    add_bullets(s, 5.3, 2.2, 4.1, 2.5, [
        "MAE: pretrain thermal-native ViT from scratch",
        "AnoDDPM: diffusion-based anomaly detection; pixel-level maps",
        "Reverse Distillation: student-teacher feature divergence",
        "All three need real GPU memory for training",
        "Unlocks with incoming NVIDIA RTX 6000 ADA 5000 (48GB)",
    ], font_size=10)
    add_textbox(s, 5.3, 4.75, 4.1, 0.35, "Blocked on hardware availability",
                font_size=10, bold=True, color=ORANGE)
    return s


def slide_next_steps(prs):
    s = blank_slide(prs)
    add_title(s, "Next Steps & Asks")
    add_subtitle(s, "What we need to unblock the next chunk of work")

    # Top: Asks (prominent callout)
    add_rect(s, 0.4, 1.3, 9.2, 1.65, fill=BG_ASK, line=ORANGE)
    add_textbox(s, 0.6, 1.45, 9.0, 0.4, "Asks : need green-light",
                font_size=17, bold=True, color=ORANGE)
    add_bullets(s, 0.6, 1.9, 9.0, 1.0, [
        ("SCADA outreach : utility-specific alarm setpoints (warning/alarm/critical tiers)", DARK),
        ("Matias + GUC engineering connection : tank-surface setpoints + deployed PV module inventory", DARK),
    ], font_size=12)

    # Bottom: engineering work
    add_textbox(s, 0.4, 3.1, 9.2, 0.4, "Engineering work (underway / queued)",
                font_size=14, bold=True, color=BLUE)

    # Two columns
    add_bullets(s, 0.4, 3.55, 4.5, 1.6, [
        "Augmentation pipeline (starts tonight; critical for transformer specialist)",
        "Source additional transformer datasets (Matias)",
        "Retrain specialist autoencoders once data + augmentation land",
    ], font_size=11)
    add_bullets(s, 5.1, 3.55, 4.5, 1.6, [
        "Ensemble pipeline MVP (stub demo tonight)",
        "PatchCore baseline (Track C v1)",
        "Track B pivot decision (Monday)",
    ], font_size=11)
    return s


def slide_thanks(prs):
    s = blank_slide(prs)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0),
                             Inches(SLIDE_W_IN), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(s, 0.5, 1.1, 9.0, 0.9, "Thank You",
                font_size=48, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 2.4, 9.0, 0.5, "Questions?",
                font_size=22, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 4.8, 9.0, 0.4,
                "HERIS Lab  |  Thermal Fault Detection System  |  April 2026",
                font_size=12, color=MID, align=PP_ALIGN.CENTER)
    return s


def main():
    if not TEMPLATE.exists():
        print(f"ERROR: template not found: {TEMPLATE}")
        return
    prs = Presentation(TEMPLATE)
    print(f"Loaded template: {TEMPLATE.name}")
    clear_slides(prs)
    print("Cleared template slides")

    builders = [
        slide_title,
        slide_recap,
        slide_specialist_models,
        slide_ensemble,
        slide_thresholds,
        slide_track_c,
        slide_next_steps,
        slide_thanks,
    ]
    for i, b in enumerate(builders, 1):
        b(prs)
        print(f"  Built slide {i}: {b.__name__}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
