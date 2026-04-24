"""Build the 2026-04-24 PV delivery presentation.

Short, focused deck (8 slides) for the live PV demo. Inherits the main
progress deck's slide master so fonts and theme match, then rebuilds
content slides from scratch.

Slides:
  1. Title
  2. What we shipped
  3. Architecture: ensemble pipeline
  4. Performance (vs. v2 baseline)
  5. Inference flow (8 steps)
  6. Demo preview (4-panel arc)
  7. Caveats + what's next
  8. Thank You / Q&A
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "docs" / "presentations" / "HERISLAB_Progress_Update.pptx"
OUT = ROOT / "docs" / "presentations" / "HERISLAB_PV_Delivery_2026-04-24.pptx"

# Theme colors (match v2 + status deck)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MID = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xA6, 0xA6, 0xA6)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
BG_CARD = RGBColor(0xF4, 0xF6, 0xF8)
BG_ASK = RGBColor(0xFF, 0xF4, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W_IN = 10.0
SLIDE_H_IN = 5.625


def clear_slides(prs):
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


def add_textbox(slide, left, top, width, height, text, *,
                font_size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = [text] if isinstance(text, str) else text
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align

        if isinstance(line, tuple):
            txt, sz, bd, clr = line
        else:
            txt, sz, bd, clr = line, font_size, bold, color

        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(sz)
        run.font.bold = bd
        if clr is not None:
            run.font.color.rgb = clr
    return tb


def add_bullets(slide, left, top, width, height, items, *, font_size=12, color=DARK,
                bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            txt, clr = item
        else:
            txt, clr = item, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.color.rgb = bullet_color
        bullet.font.size = Pt(font_size)
        bullet.font.bold = True
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.color.rgb = clr
        run.font.name = "Calibri"
    return tb


def add_numbered(slide, left, top, width, height, items, *, font_size=11, color=DARK,
                 number_color=BLUE):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(5)
        num = p.add_run()
        num.text = f"{i+1}.  "
        num.font.color.rgb = number_color
        num.font.size = Pt(font_size)
        num.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_rect(slide, left, top, width, height, fill=BG_CARD, line=None):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_title(slide, text, color=BLUE, size=28):
    return add_textbox(slide, 0.4, 0.25, 9.2, 0.5, text,
                       font_size=size, bold=True, color=color)


def add_subtitle(slide, text, color=MID, size=13):
    return add_textbox(slide, 0.4, 0.8, 9.2, 0.35, text,
                       font_size=size, bold=False, color=color)


# ---- Slides ----

def slide_title(prs):
    s = blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0),
                             Inches(SLIDE_W_IN), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(s, 0.5, 1.0, 9.0, 0.9, "PV Fault Detection: Delivery",
                font_size=42, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 1.7, 9.0, 0.5, "HERIS Lab  |  Thermal Fault Detection System",
                font_size=20, bold=False, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 2.4, 9.0, 0.4, "April 24, 2026",
                font_size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 4.6, 9.0, 0.4,
                "Ryan Balungeli  |  Matias  |  Pablo",
                font_size=14, color=MID, align=PP_ALIGN.CENTER)
    return s


def slide_shipped(prs):
    s = blank_slide(prs)
    add_title(s, "What We Built")
    add_subtitle(s, "A working PV fault detection pipeline, from image to decision")

    # Four cards
    card_y = 1.4
    card_h = 3.6
    cards = [
        {
            "title": "PV-Only Model",
            "header_color": BLUE,
            "body": [
                "Trained only on PV panel images",
                "8,019 training images",
                "Trained in 10.5 minutes on our GPU",
                "Separate from the earlier all-equipment model",
            ],
        },
        {
            "title": "Two-Layer Decision",
            "header_color": ACCENT,
            "body": [
                "The model's anomaly check",
                "A temperature-based rule check",
                "Combined answer with an explanation",
                "Either layer can flag a problem",
            ],
        },
        {
            "title": "Tested on Real Data",
            "header_color": GREEN,
            "body": [
                "892 healthy panels held back for testing",
                "1,000 real fault images from an open dataset",
                "Faults include cracks, hot spots, and shading",
                "Model had never seen any of them",
            ],
        },
        {
            "title": "Live Demo",
            "header_color": ORANGE,
            "body": [
                "One command runs the whole pipeline",
                "Shows 4 example images",
                "Produces a clear visual per image",
                "Can also run on any single image",
            ],
        },
    ]

    card_w = (SLIDE_W_IN - 0.4 - 0.4 - 0.15 * 3) / 4
    x = 0.4
    for c in cards:
        add_rect(s, x, card_y, card_w, card_h, fill=BG_CARD)
        add_textbox(s, x + 0.15, card_y + 0.15, card_w - 0.3, 0.4,
                    c["title"], font_size=14, bold=True, color=c["header_color"])
        add_bullets(s, x + 0.15, card_y + 0.7, card_w - 0.3, card_h - 0.8,
                    c["body"], font_size=10)
        x += card_w + 0.15
    return s


def slide_architecture(prs):
    s = blank_slide(prs)
    add_title(s, "Architecture: Ensemble Pipeline")
    add_subtitle(s, "The model catches unusual patterns; the rules catch clear temperature limits. They cover each other.")

    # Flow diagram (horizontal)
    box_h = 0.75
    box_y = 1.5

    # Input
    add_rect(s, 0.3, box_y, 1.6, box_h, fill=ACCENT)
    add_textbox(s, 0.3, box_y + 0.1, 1.6, 0.4, "Raw Input",
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.3, box_y + 0.45, 1.6, 0.3, "Thermal image",
                font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(s, 1.95, box_y + 0.2, 0.3, 0.4, ">",
                font_size=24, bold=True, color=MID, align=PP_ALIGN.CENTER)

    # ML path (top)
    add_rect(s, 2.3, 1.05, 2.2, 0.6, fill=GREEN)
    add_textbox(s, 2.3, 1.15, 2.2, 0.4, "Model Check",
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 2.3, 1.75, 2.2, 0.35, "How well it can rebuild the image",
                font_size=9, color=DARK, align=PP_ALIGN.CENTER)

    # Rules path (bottom)
    add_rect(s, 2.3, 2.35, 2.2, 0.6, fill=ORANGE)
    add_textbox(s, 2.3, 2.45, 2.2, 0.4, "Temperature Rule",
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 2.3, 3.05, 2.2, 0.35, "Is any spot too hot vs limits",
                font_size=9, color=DARK, align=PP_ALIGN.CENTER)

    # Arrows in
    add_textbox(s, 4.55, 1.15, 0.3, 0.4, ">", font_size=20, bold=True, color=MID)
    add_textbox(s, 4.55, 2.45, 0.3, 0.4, ">", font_size=20, bold=True, color=MID)

    # Ensemble
    add_rect(s, 4.95, box_y, 1.8, box_h, fill=BLUE)
    add_textbox(s, 4.95, box_y + 0.15, 1.8, 0.4, "Ensemble",
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, 4.95, box_y + 0.5, 1.8, 0.3, "Combined decision",
                font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow out
    add_textbox(s, 6.8, box_y + 0.2, 0.3, 0.4, ">", font_size=20, bold=True, color=MID)

    # Output
    add_rect(s, 7.15, box_y, 2.5, box_h, fill=BG_CARD, line=DARK)
    add_textbox(s, 7.15, box_y + 0.1, 2.5, 0.4, "Output",
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, 7.15, box_y + 0.45, 2.5, 0.3, "Verdict, severity, explanation",
                font_size=9, color=MID, align=PP_ALIGN.CENTER)

    # Bottom: why both
    add_rect(s, 0.4, 3.7, 9.2, 1.7, fill=BG_CARD)
    add_textbox(s, 0.6, 3.85, 9.0, 0.4, "Why both layers",
                font_size=14, bold=True, color=BLUE)
    add_textbox(s, 0.6, 4.25, 4.4, 0.3, "The model catches:",
                font_size=11, bold=True, color=GREEN)
    add_bullets(s, 0.6, 4.5, 4.4, 0.9, [
        "Unusual visual patterns (cracks, damage, shading)",
        "Problems the rules might not have a number for",
        "Early warning signs before anything gets really hot",
    ], font_size=9)
    add_textbox(s, 5.1, 4.25, 4.4, 0.3, "The rules catch:",
                font_size=11, bold=True, color=ORANGE)
    add_bullets(s, 5.1, 4.5, 4.4, 0.9, [
        "Clear overheat cases with known temperature limits",
        "Easy-to-explain decisions for a technician or auditor",
        "Hard safety limits that should never be crossed",
    ], font_size=9)
    return s


def slide_performance(prs):
    s = blank_slide(prs)
    add_title(s, "How Well It Works")
    add_subtitle(s, "Tested on 892 healthy panels set aside for testing plus 1,000 real fault images")

    # Left: metrics comparison table
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Scores (PV-only model)",
                font_size=15, bold=True, color=BLUE)

    # Header row
    add_textbox(s, 0.7, 2.0, 1.6, 0.3, "Metric", font_size=11, bold=True, color=DARK)
    add_textbox(s, 2.3, 2.0, 1.3, 0.3, "Old model", font_size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, 3.6, 2.0, 1.3, 0.3, "PV-only model", font_size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    rows = [
        ("AUROC", "0.9023", "0.9736", GREEN),
        ("F1", "0.9807", "0.9441", MID),
        ("Precision", "0.9621", "0.9178", MID),
        ("Recall", "1.0000", "0.9720", MID),
    ]
    y = 2.5
    for label, v2, pv, hl in rows:
        add_textbox(s, 0.7, y, 1.6, 0.3, label, font_size=11, color=DARK)
        add_textbox(s, 2.3, y, 1.3, 0.3, v2, font_size=11, color=MID, align=PP_ALIGN.CENTER)
        add_textbox(s, 3.6, y, 1.3, 0.3, pv, font_size=11, bold=True, color=hl, align=PP_ALIGN.CENTER)
        y += 0.4

    add_textbox(s, 0.6, 4.25, 4.1, 0.9,
                "Healthy panels produce low error (0.015 average). Faulty panels produce "
                "much higher error (0.147 average). About a 10x gap between the two groups.",
                font_size=10, color=MID)

    # Right: confusion matrix
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "Confusion matrix",
                font_size=15, bold=True, color=BLUE)
    add_textbox(s, 5.3, 1.85, 4.1, 0.3, "Using the best threshold for balanced accuracy (0.028)",
                font_size=10, color=MID)

    # 2x2 grid
    cell_w, cell_h = 1.7, 0.85
    cx, cy = 6.2, 2.4

    # Corner labels
    add_textbox(s, 5.3, cy + cell_h * 0.3, 0.9, 0.3, "Actual", font_size=10, bold=True, color=MID)
    add_textbox(s, 5.3, cy + cell_h * 1.3, 0.9, 0.3, "Actual", font_size=10, bold=True, color=MID)
    add_textbox(s, 5.3, cy + cell_h * 0.5, 0.9, 0.3, "Fault", font_size=10, color=DARK)
    add_textbox(s, 5.3, cy + cell_h * 1.5, 0.9, 0.3, "Normal", font_size=10, color=DARK)

    # Column headers
    add_textbox(s, cx, cy - 0.45, cell_w, 0.25, "Pred Fault", font_size=10, bold=True, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + cell_w, cy - 0.45, cell_w, 0.25, "Pred Normal", font_size=10, bold=True, color=MID, align=PP_ALIGN.CENTER)

    # 4 cells with numbers
    cells = [
        (cx, cy, "TP", "972", GREEN),
        (cx + cell_w, cy, "FN", "28", RED),
        (cx, cy + cell_h, "FP", "87", ORANGE),
        (cx + cell_w, cy + cell_h, "TN", "805", GREEN),
    ]
    for x, y_, lbl, num, col in cells:
        add_rect(s, x, y_, cell_w, cell_h, fill=WHITE, line=MID)
        add_textbox(s, x, y_ + 0.1, cell_w, 0.3, lbl, font_size=10, bold=True, color=MID, align=PP_ALIGN.CENTER)
        add_textbox(s, x, y_ + 0.4, cell_w, 0.4, num, font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

    add_textbox(s, 5.3, 4.7, 4.1, 0.4,
                "97 percent of faults caught;  90 percent of normals correctly passed.",
                font_size=10, color=MID)

    # Source attribution footer
    add_textbox(s, 0.4, 5.3, 9.2, 0.3,
                "Test faults: PVMD open dataset (Mendeley, 2024). See References slide.",
                font_size=9, color=LIGHT, align=PP_ALIGN.CENTER)
    return s


def slide_inference(prs):
    s = blank_slide(prs)
    add_title(s, "How Each Image Is Processed")
    add_subtitle(s, "The steps the code runs, per image, from file to final answer")

    # Left: numbered steps 1-5 (prep + model)
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Prepare the image",
                font_size=14, bold=True, color=BLUE)
    add_numbered(s, 0.6, 1.95, 4.1, 3.1, [
        "Load the image (from camera or file)",
        "Strip color: keep only brightness per pixel",
        "Standardize the brightness range to 0 to 255",
        "Resize to the model's expected size; balance the numbers",
        "Run the image through the model",
    ], font_size=11)

    # Right: steps 6-9 (decision)
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "Make the decision",
                font_size=14, bold=True, color=BLUE)
    add_numbered(s, 5.3, 1.95, 4.1, 3.1, [
        "Measure how well the model rebuilt the image",
        "Model answer: high error means fault (vs. a calibrated cutoff)",
        "Rule answer: does any spot exceed temperature limits?",
        "Combine both answers into a final verdict",
    ], font_size=11)

    # Bottom note
    add_textbox(s, 0.4, 5.25, 9.2, 0.3,
                "Under a second per image on our GPU. The demo runs this on 4 samples and draws a picture showing each step.",
                font_size=10, color=MID, align=PP_ALIGN.CENTER)
    return s


def slide_demo(prs):
    s = blank_slide(prs)
    add_title(s, "Live Demo: 4 Example Images")
    add_subtitle(s, "Each one shows what the model and rules each say, and how they combine")

    # Four cards, one per panel
    cards = [
        {
            "tag": "[1]",
            "title": "Healthy panel",
            "ml": "NORMAL",
            "ml_color": GREEN,
            "rule": "NORMAL",
            "rule_color": GREEN,
            "ensemble": "NORMAL",
            "ensemble_color": GREEN,
            "note": "Both layers agree: nothing wrong",
        },
        {
            "tag": "[2]",
            "title": "Healthy but unusual",
            "ml": "FAULT",
            "ml_color": RED,
            "rule": "NORMAL",
            "rule_color": GREEN,
            "ensemble": "Flagged by model",
            "ensemble_color": ORANGE,
            "note": "Model noticed something odd; would send for a check before any crew is dispatched",
        },
        {
            "tag": "[3]",
            "title": "Mild hot spot",
            "ml": "FAULT",
            "ml_color": RED,
            "rule": "WARNING",
            "rule_color": ORANGE,
            "ensemble": "Both agree",
            "ensemble_color": RED,
            "note": "Both layers flag: high-confidence fault",
        },
        {
            "tag": "[4]",
            "title": "Severe crack",
            "ml": "FAULT",
            "ml_color": RED,
            "rule": "WARNING",
            "rule_color": ORANGE,
            "ensemble": "Model catches it",
            "ensemble_color": RED,
            "note": "Cracks look dim, not hot, so the rules alone would miss this. The model catches it.",
        },
    ]

    card_w = (SLIDE_W_IN - 0.4 - 0.4 - 0.15 * 3) / 4
    card_h = 3.6
    card_y = 1.4
    x = 0.4

    for c in cards:
        add_rect(s, x, card_y, card_w, card_h, fill=BG_CARD)
        add_textbox(s, x + 0.15, card_y + 0.15, card_w - 0.3, 0.3,
                    c["tag"], font_size=11, bold=True, color=MID)
        add_textbox(s, x + 0.15, card_y + 0.45, card_w - 0.3, 0.35,
                    c["title"], font_size=12, bold=True, color=BLUE)

        # Verdicts stack
        vy = card_y + 1.0
        add_textbox(s, x + 0.15, vy, card_w - 0.3, 0.25, "ML", font_size=9, color=MID)
        add_textbox(s, x + 0.15, vy + 0.25, card_w - 0.3, 0.3, c["ml"],
                    font_size=13, bold=True, color=c["ml_color"])
        add_textbox(s, x + 0.15, vy + 0.65, card_w - 0.3, 0.25, "Rule", font_size=9, color=MID)
        add_textbox(s, x + 0.15, vy + 0.9, card_w - 0.3, 0.3, c["rule"],
                    font_size=13, bold=True, color=c["rule_color"])
        add_textbox(s, x + 0.15, vy + 1.3, card_w - 0.3, 0.25, "Ensemble", font_size=9, color=MID)
        add_textbox(s, x + 0.15, vy + 1.55, card_w - 0.3, 0.3, c["ensemble"],
                    font_size=13, bold=True, color=c["ensemble_color"])

        add_textbox(s, x + 0.15, card_y + card_h - 0.7, card_w - 0.3, 0.6,
                    c["note"], font_size=9, color=MID)
        x += card_w + 0.15

    add_textbox(s, 0.4, 5.25, 9.2, 0.3,
                "Running the demo now.",
                font_size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    return s


def slide_caveats(prs):
    s = blank_slide(prs)
    add_title(s, "What's Rough + What's Next")
    add_subtitle(s, "Honest about shortcuts taken today and what's immediately next")

    # Left: rough edges (orange)
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_ASK, line=ORANGE)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Rough edges to fix",
                font_size=15, bold=True, color=ORANGE)
    add_bullets(s, 0.6, 1.9, 4.1, 3.1, [
        "The temperature rule uses an estimate, not real degrees. The open "
        "dataset we used doesn't include real temperature readings. A real "
        "deployment camera provides them directly.",
        "You have to tell the system what equipment the image is (PV or "
        "transformer). A real deployment would detect that automatically.",
        "There's no web service yet. The demo runs on a laptop. A real "
        "deployment needs a server behind an API.",
    ], font_size=10)

    # Right: next up (green) — ordered by critical path
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "Next up (in order)",
                font_size=15, bold=True, color=GREEN)
    add_bullets(s, 5.3, 1.9, 4.1, 3.1, [
        "Build the data augmentation pipeline. Still pending from last "
        "round. The transformer model has only 18 training images and "
        "can't train without this.",
        "Keep sourcing more transformer images (Matias leading).",
        "Train the transformer-only model once the above two are ready.",
        "Tune the alert cutoff based on whether the utility prefers "
        "'catch everything' or 'don't cry wolf'.",
        "Wrap the pipeline in a web service for real-time use.",
    ], font_size=10)
    return s


def slide_references(prs):
    s = blank_slide(prs)
    add_title(s, "References and Credits", size=26)
    add_subtitle(s, "External data sources and standards used in this delivery")

    # Left column: datasets
    add_rect(s, 0.4, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 0.6, 1.45, 4.1, 0.4, "Datasets",
                font_size=14, bold=True, color=BLUE)

    # PVMD
    add_textbox(s, 0.6, 1.9, 4.1, 0.3, "PVMD: Photovoltaic Module Dataset",
                font_size=11, bold=True, color=DARK)
    add_textbox(s, 0.6, 2.2, 4.1, 0.5,
                "Mendeley Data, 2024. 1,000 labeled fault images (cracks, hot spots, and shadings). Used to test this model.",
                font_size=9, color=MID)
    add_textbox(s, 0.6, 2.75, 4.1, 0.25,
                "data.mendeley.com/datasets/5ssmfpgrpc/1",
                font_size=8, color=ACCENT)

    # PV O&M
    add_textbox(s, 0.6, 3.15, 4.1, 0.3, "PV System O&M Inspection",
                font_size=11, bold=True, color=DARK)
    add_textbox(s, 0.6, 3.45, 4.1, 0.45,
                "7,836 healthy PV thermal images used to train the model.",
                font_size=9, color=MID)

    # PV Thermal
    add_textbox(s, 0.6, 3.95, 4.1, 0.3, "PV System Thermal Inspection",
                font_size=11, bold=True, color=DARK)
    add_textbox(s, 0.6, 4.25, 4.1, 0.45,
                "1,075 healthy PV thermal images used to train the model.",
                font_size=9, color=MID)

    # Right column: standards and threshold sources
    add_rect(s, 5.1, 1.3, 4.5, 3.9, fill=BG_CARD)
    add_textbox(s, 5.3, 1.45, 4.1, 0.4, "Rule-layer threshold sources",
                font_size=14, bold=True, color=BLUE)

    add_textbox(s, 5.3, 1.9, 4.1, 0.3, "Standards",
                font_size=11, bold=True, color=DARK)
    add_bullets(s, 5.3, 2.2, 4.1, 1.6, [
        "IEC 61215: PV module design qualification",
        "IEC 61730: PV module safety qualification",
        "IEC TS 63126:2025: operation at elevated temperatures",
        "UL 61730: PV module safety standard (US)",
    ], font_size=9)

    add_textbox(s, 5.3, 3.85, 4.1, 0.3, "Manufacturer datasheets",
                font_size=11, bold=True, color=DARK)
    add_bullets(s, 5.3, 4.15, 4.1, 1.0, [
        "LG NeON R, Canadian Solar HiKu, JinkoSolar Tiger Neo, "
        "Trina Vertex, First Solar Series 6",
    ], font_size=9)

    return s


def slide_thanks(prs):
    s = blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0),
                             Inches(SLIDE_W_IN), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(s, 0.5, 1.1, 9.0, 0.9, "Thank You",
                font_size=48, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 2.4, 9.0, 0.5, "Questions?",
                font_size=22, color=MID, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 4.8, 9.0, 0.4,
                "HERIS Lab  |  Thermal Fault Detection System  |  April 2026",
                font_size=12, color=MID, align=PP_ALIGN.CENTER)
    return s


def main():
    if not TEMPLATE.exists():
        print(f"ERROR: template not found: {TEMPLATE}")
        return
    prs = Presentation(TEMPLATE)
    print(f"Loaded template: {TEMPLATE.name}")
    clear_slides(prs)
    print("Cleared template slides")

    builders = [
        slide_title,
        slide_shipped,
        slide_architecture,
        slide_performance,
        slide_inference,
        slide_demo,
        slide_caveats,
        slide_references,
        slide_thanks,
    ]
    for i, b in enumerate(builders, 1):
        b(prs)
        print(f"  Built slide {i}: {b.__name__}")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
