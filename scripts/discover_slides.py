"""Discovery pass: dump all text per slide for the existing pptx.

Used to identify exact strings to replace in update_slides_v2.py.
Writes to results/slide_discovery.txt (utf-8 safe).
"""

from pptx import Presentation
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "docs" / "HERISLAB_Progress_Update.pptx"
OUT = ROOT / "results" / "slide_discovery.txt"

OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation(PPTX)
lines = [f"Total slides: {len(prs.slides)}", ""]

for i, slide in enumerate(prs.slides, 1):
    lines.append(f"===== SLIDE {i} =====")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(f"  [shape {shape_idx} text] {text!r}")
        if shape.has_table:
            for row_idx, row in enumerate(shape.table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                lines.append(f"  [shape {shape_idx} table row {row_idx}] {cells!r}")
    lines.append("")

OUT.write_text("\n".join(lines), encoding="utf-8")
print(f"Wrote: {OUT}")
print(f"Total slides: {len(prs.slides)}")
